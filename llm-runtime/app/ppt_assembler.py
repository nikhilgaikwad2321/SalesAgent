import os
import io
from datetime import datetime
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt


def create_charts(slide_data: dict) -> List[str]:
    """
    Generate deterministic charts using matplotlib.
    Returns list of file paths to generated chart images.
    """
    chart_paths = []
    
    # Extract chart data from slides
    coverage_data = None
    for slide in slide_data.get("slides", []):
        if slide.get("data") and "coverage" in slide.get("title", "").lower():
            coverage_data = slide.get("data")
            break
    
    # Create output directory
    os.makedirs("temp/charts", exist_ok=True)
    
    # Chart 1: Coverage vs Premium (if data available)
    if coverage_data and "coverage_amounts" in coverage_data:
        fig, ax = plt.subplots(figsize=(8, 5))
        coverage_amounts = coverage_data.get("coverage_amounts", [10, 25, 50, 100])
        premiums = coverage_data.get("premiums", [500, 1000, 1800, 3200])
        
        ax.bar(range(len(coverage_amounts)), premiums, color='#2563eb', alpha=0.8)
        ax.set_xlabel('Coverage Amount (Lakhs)', fontsize=12, fontweight='bold')
        ax.set_ylabel('Annual Premium (₹)', fontsize=12, fontweight='bold')
        ax.set_title('Coverage vs Premium Illustration', fontsize=14, fontweight='bold', pad=20)
        ax.set_xticks(range(len(coverage_amounts)))
        ax.set_xticklabels([f'₹{x}L' for x in coverage_amounts])
        ax.grid(axis='y', alpha=0.3)
        
        chart_path = f"temp/charts/coverage_premium_{datetime.now().strftime('%Y%m%d%H%M%S')}.png"
        plt.tight_layout()
        plt.savefig(chart_path, dpi=150, bbox_inches='tight')
        plt.close()
        chart_paths.append(chart_path)
    else:
        # Default illustrative chart
        fig, ax = plt.subplots(figsize=(8, 5))
        coverage_amounts = [10, 25, 50, 100]
        premiums = [500, 1000, 1800, 3200]
        
        ax.bar(range(len(coverage_amounts)), premiums, color='#2563eb', alpha=0.8)
        ax.set_xlabel('Coverage Amount (Lakhs)', fontsize=12, fontweight='bold')
        ax.set_ylabel('Indicative Annual Premium (₹)', fontsize=12, fontweight='bold')
        ax.set_title('Coverage vs Premium - Illustrative', fontsize=14, fontweight='bold', pad=20)
        ax.set_xticks(range(len(coverage_amounts)))
        ax.set_xticklabels([f'₹{x}L' for x in coverage_amounts])
        ax.grid(axis='y', alpha=0.3)
        
        chart_path = f"temp/charts/coverage_premium_{datetime.now().strftime('%Y%m%d%H%M%S')}.png"
        plt.tight_layout()
        plt.savefig(chart_path, dpi=150, bbox_inches='tight')
        plt.close()
        chart_paths.append(chart_path)
    
    # Chart 2: Benefits/Allocation Breakdown (Donut Chart)
    if True: # Always generate for demo
        fig, ax = plt.subplots(figsize=(8, 6))
        benefits = ['Death Benefit', 'Critical Illness', 'Accidental Death', 'Disability']
        sizes = [50, 25, 15, 10]
        colors = ['#2563eb', '#3b82f6', '#60a5fa', '#93c5fd']
        
        # Donut Chart
        wedges, texts, autotexts = ax.pie(sizes, labels=benefits, colors=colors, autopct='%1.1f%%',
               shadow=False, startangle=90, pctdistance=0.85, textprops={'fontsize': 11, 'fontweight': 'bold'})
        
        # Draw white circle in center
        centre_circle = plt.Circle((0,0),0.70,fc='white')
        fig.gca().add_artist(centre_circle)
        
        ax.axis('equal')  
        ax.set_title('Benefit Allocation', fontsize=14, fontweight='bold', pad=20)
        
        chart_path = f"temp/charts/benefits_{datetime.now().strftime('%Y%m%d%H%M%S')}.png"
        plt.tight_layout()
        plt.savefig(chart_path, dpi=150, bbox_inches='tight')
        plt.close()
        chart_paths.append(chart_path)

    # Chart 3: Returns Comparison (Dual Line Graph 4% vs 8%)
    if True:
        fig, ax = plt.subplots(figsize=(8, 5))
        years = range(1, 21) # 20 year view
        returns_4 = [100 * (1.04 ** i) for i in years] # Proxy logic
        returns_8 = [100 * (1.08 ** i) for i in years] # Proxy logic
        
        ax.plot(years, returns_8, marker='', color='#16a34a', linewidth=3, label='Optimistic (8%)')
        ax.plot(years, returns_4, marker='', color='#2563eb', linewidth=3, linestyle='--', label='Conservative (4%)')
        
        ax.set_xlabel('Policy Year', fontsize=12, fontweight='bold')
        ax.set_ylabel('Fund Value (Index)', fontsize=12, fontweight='bold')
        ax.set_title('Projected Returns Growth (4% vs 8%)', fontsize=14, fontweight='bold', pad=20)
        ax.legend(fontsize=11)
        ax.grid(True, alpha=0.3)
        
        chart_path = f"temp/charts/returns_comparison_{datetime.now().strftime('%Y%m%d%H%M%S')}.png"
        plt.tight_layout()
        plt.savefig(chart_path, dpi=150, bbox_inches='tight')
        plt.close()
        chart_paths.append(chart_path)

    # Chart 4: Premium Payment Timeline (Horizontal Bar)
    if True:
        fig, ax = plt.subplots(figsize=(8, 3))
        stages = ['Pay Premium', 'Stay Protected', 'Maturity Benefit']
        start = [0, 0, 15]
        duration = [10, 20, 0.5] # Pay for 10, Cover for 20 (Simulated)
        colors = ['#f59e0b', '#2563eb', '#16a34a']
        
        ax.barh(stages, duration, left=start, color=colors, height=0.5)
        ax.set_xlabel('Years', fontsize=12, fontweight='bold')
        ax.set_title('Premium Payment & Cover Timeline', fontsize=14, fontweight='bold', pad=10)
        ax.grid(axis='x', alpha=0.3)
        
        chart_path = f"temp/charts/timeline_{datetime.now().strftime('%Y%m%d%H%M%S')}.png"
        plt.tight_layout()
        plt.savefig(chart_path, dpi=150, bbox_inches='tight')
        plt.close()
        chart_paths.append(chart_path)

    return chart_paths


def assemble_ppt(slide_content: dict, chart_paths: List[str]) -> str:
    """
    Assemble the PowerPoint presentation using python-pptx.
    Returns the file path to the generated PPT.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    PRIMARY_COLOR = RGBColor(37, 99, 235)  # Blue
    TEXT_COLOR = RGBColor(31, 41, 55)  # Dark gray
    
    client_name = slide_content.get("client_name", "Valued Client")
    policy_name = slide_content.get("policy_name", "Insurance Policy")
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    background = slide.shapes.add_shape(1, left, top, width, height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = policy_name
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Prepared for: {client_name}"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = TEXT_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%B %d, %Y")
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = RGBColor(107, 114, 128)
    date_para.alignment = PP_ALIGN.CENTER
    
    # Process content slides
    for slide_data in slide_content.get("slides", []):
        slide_layout = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Clear default placeholders
        for shape in slide.shapes:
            if shape.is_placeholder:
                sp = shape.element
                sp.getparent().remove(sp)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get("title", "")
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        
        title_para.font.color.rgb = PRIMARY_COLOR
        
        # Add content based on type
        content_type = slide_data.get("content_type", "text")
        content = slide_data.get("content", "")
        
        if content_type == "kpi_cards":
            # Draw KPI Cards (Rounded Rectangles)
            # Content is expected to be a list of dicts: [{"label": "Sum Assured", "value": "₹50 Lakh"}, ...]
            kpi_data = content if isinstance(content, list) else []
            card_width = Inches(2.5)
            card_height = Inches(1.5)
            gap = Inches(0.5)
            start_x = (prs.slide_width - ((card_width * len(kpi_data)) + (gap * (len(kpi_data) - 1)))) / 2
            start_y = Inches(2.5)
            
            for i, kpi in enumerate(kpi_data):
                left_pos = start_x + (i * (card_width + gap))
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, start_y, card_width, card_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(239, 246, 255) # Light blue bg
                shape.line.color.rgb = PRIMARY_COLOR
                shape.line.width = Pt(1.5)
                
                # Value (Top/Bold)
                text_frame = shape.text_frame
                text_frame.text = kpi.get("value", "")
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_COLOR
                
                # Label (Bottom)
                p2 = text_frame.add_paragraph()
                p2.text = kpi.get("label", "")
                p2.alignment = PP_ALIGN.CENTER
                p2.font.size = Pt(14)
                p2.font.color.rgb = TEXT_COLOR
                
        elif content_type == "bullet_points" and isinstance(content, list):
            content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(4.5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, bullet in enumerate(content):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_COLOR
                p.space_before = Pt(12)
        else:
            content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(4.5))
            text_frame = content_box.text_frame
            text_frame.text = str(content)
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
    
    # Add chart slides
    if chart_paths:
        # We will add slides for charts systematically
        # 1. Coverage vs Premium (from original code, if exists[0])
        # 2. Donut (Allocation)
        # 3. Dual Line (Returns)
        # 4. Timeline
        
        chart_titles = [
            "Coverage vs Premium Analysis",
            "Your Investment Allocation",
            "Projected Returns Simulation",
            "Long-Term Protection Timeline"
        ]
        
        for i, chart_path in enumerate(chart_paths):
            if i >= len(chart_titles): break # Safety check
            
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            
            # Clear placeholders
            for shape in slide.shapes:
                if shape.is_placeholder:
                    sp = shape.element
                    sp.getparent().remove(sp)
            
            # Add title
            title_text = chart_titles[i]
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = PRIMARY_COLOR
            
            # Add chart image
            if os.path.exists(chart_path):
                # Center the image
                pic = slide.shapes.add_picture(chart_path, Inches(1), Inches(1.8), width=Inches(8))
                
                # Add a caption/explanation below chart
                caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
                caption_frame = caption_box.text_frame
                
                captions = [
                    "Higher coverage ensures your family is safe, with affordable premiums.",
                    "Your money is diversified to ensure both safety and growth.",
                    "See how your fund grows significantly over 20 years with compound interest.",
                    "You pay for a limited period, but get protection for the full term."
                ]
                caption_frame.text = captions[i] if i < len(captions) else ""
                cp = caption_frame.paragraphs[0]
                cp.font.size = Pt(16)
                cp.font.italic = True
                cp.alignment = PP_ALIGN.CENTER
                cp.font.color.rgb = TEXT_COLOR
    
    # Disclaimer Slide
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    for shape in slide.shapes:
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Important Disclaimer"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    
    disclaimer_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    disclaimer_frame = disclaimer_box.text_frame
    disclaimer_frame.word_wrap = True
    
    disclaimer_text = """All illustrations, premium amounts, and coverage details presented in this document are indicative and based on approved policy information available at the time of preparation.

Actual premium rates, coverage amounts, and policy terms may vary based on individual circumstances, age, health conditions, and underwriting guidelines.

This presentation does not constitute a guarantee of coverage or returns. Please refer to the official policy document for complete terms and conditions.

For accurate quotations and personalized policy recommendations, please consult with our licensed insurance advisors."""
    
    disclaimer_frame.text = disclaimer_text
    p = disclaimer_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(107, 114, 128)
    p.line_spacing = 1.5
    
    # Save presentation
    os.makedirs("temp/ppts", exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    safe_client_name = "".join(c for c in client_name if c.isalnum() or c in (' ', '-', '_')).strip()
    safe_client_name = safe_client_name.replace(' ', '_')
    filename = f"Insurance_Presentation_{safe_client_name}_{timestamp}.pptx"
    file_path = os.path.join("temp/ppts", filename)
    
    prs.save(file_path)
    
    return file_path
